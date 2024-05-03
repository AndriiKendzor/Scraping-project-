from selenium import webdriver
from chromedriver_py import binary_path
from selenium.webdriver.common.by import By
import time

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

class DataScraper:
    def __init__(self, url):
        self.url = url
        self.title_top = []
        self.titles = []
        self.informs = []
        self.top_title = ""
        self.titless = []
        self.informss = []

    def scrape_data(self):
        svc = webdriver.ChromeService(executable_path=binary_path)
        driver = webdriver.Chrome(service=svc)
        driver.get(self.url)

        self.title_top = driver.find_elements(By.XPATH, '//div[@class="lms-header-title"]')
        self.titles = driver.find_elements(By.XPATH, '//h2')
        self.informs = driver.find_elements(By.XPATH, '//p')

        time.sleep(5)

        for title in self.title_top:
            self.top_title = title.text

        for title in self.titles:
            self.titless.append(title.text)

        for info in self.informs:
            self.informss.append(info.text)

        driver.close()


class PresentationCreator:
    def __init__(self):
        self.prs = Presentation()

    def add_title_slide(self, title, title_font_size):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0]
        title_placeholder.text = title

        title_paragraph = title_placeholder.text_frame.paragraphs[0]
        for run in title_paragraph.runs:
            run.font.size = Pt(title_font_size)
            run.font.color.rgb = RGBColor(255, 213, 0)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 14, 40)

    def add_text_slide(self, text, text_font_size):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        text_placeholder = slide.placeholders[1]
        text_placeholder.text = text

        for paragraph in text_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(text_font_size)
                run.font.color.rgb = RGBColor(255, 213, 0)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 14, 40)

    def add_slide(self, title, content, title_font_size, content_font_size):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0]
        title_placeholder.text = title

        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content

        title_paragraph = title_placeholder.text_frame.paragraphs[0]
        for run in title_paragraph.runs:
            run.font.size = Pt(title_font_size)
            run.font.color.rgb = RGBColor(255, 213, 0)

        for paragraph in content_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(content_font_size)
                run.font.color.rgb = RGBColor(255, 213, 0)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 14, 40)






url = 'https://kanga.exchange/university/courses/poziom-podstawowy/lessons/3-satoshi-nakamoto-kim-jest-tworca-bitcoina/'
scraper = DataScraper(url)
scraper.scrape_data()


presentation_creator = PresentationCreator()

presentation_creator.add_title_slide(scraper.top_title[2:], title_font_size=40)
presentation_creator.add_text_slide(scraper.informss[1], text_font_size=15)

slide1_info = '\n'.join(scraper.informss[2:5])
presentation_creator.add_slide(scraper.titless[1], slide1_info, title_font_size=20, content_font_size=15)

slide2_info = '\n'.join(scraper.informss[5:9])
presentation_creator.add_slide(scraper.titless[2], slide2_info, title_font_size=20, content_font_size=14)

slide2_info2 = '\n'.join(scraper.informss[9:12])
presentation_creator.add_text_slide(slide2_info2, text_font_size=15)

slide3_info = '\n'.join(scraper.informss[12:19])
presentation_creator.add_slide(scraper.titless[3], slide3_info, title_font_size=20, content_font_size=15)

slide3_info2 = '\n'.join(scraper.informss[19:27])
presentation_creator.add_text_slide(slide3_info2, text_font_size=14)

slide4_info = '\n'.join(scraper.informss[28:31])
presentation_creator.add_slide(scraper.informss[27], slide4_info, title_font_size=20, content_font_size=15)

presentation_creator.add_slide('Źródło', 'https://kanga.exchange/university/courses/poziom-podstawowy/lessons/3-satoshi-nakamoto-kim-jest-tworca-bitcoina/', title_font_size=20, content_font_size=15)
presentation_creator.add_title_slide('Wykonane przez Andrii Kendzor', title_font_size=30)
presentation_creator.prs.save('Kto jest prawdziwym twórcą Bitcoina (63656).pptx')
print("success")

